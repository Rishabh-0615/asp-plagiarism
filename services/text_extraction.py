"""Text Extraction Service for various file formats"""
import os
import logging
from pathlib import Path
from urllib.parse import urlparse

try:
    import PyPDF2
except ImportError:
    PyPDF2 = None

try:
    from docx import Document
except ImportError:
    Document = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

logger = logging.getLogger(__name__)

class TextExtractionService:

    CODE_EXTENSIONS = {
        '.java', '.py', '.js', '.jsx', '.ts', '.tsx', '.c', '.cpp', '.cc',
        '.cxx', '.h', '.hpp', '.cs', '.go', '.rs', '.php', '.rb', '.swift',
        '.kt', '.kts', '.scala', '.sql', '.sh', '.ps1', '.bat'
    }

    @staticmethod
    def _extension_from_content_type(content_type: str) -> str:
        if not content_type:
            return ""
        ctype = content_type.lower()
        if "application/pdf" in ctype:
            return ".pdf"
        if "application/vnd.openxmlformats-officedocument.wordprocessingml.document" in ctype:
            return ".docx"
        if "application/vnd.openxmlformats-officedocument.presentationml.presentation" in ctype:
            return ".pptx"
        if "text/plain" in ctype or "text/" in ctype:
            return ".txt"
        return ""

    @staticmethod
    def _extension_from_magic(content: bytes) -> str:
        if not content:
            return ""
        if content.startswith(b"%PDF"):
            return ".pdf"
        # Office Open XML files are ZIP containers. Leave as txt fallback unless
        # content-type tells us whether it is docx or pptx.
        if content.startswith(b"PK"):
            return ".txt"
        return ".txt"
    
    @staticmethod
    def extract_from_file(file_path: str) -> str:
        """
        Extract text from various file formats
        
        Args:
            file_path (str): Path to the file
            
        Returns:
            str: Extracted text
        """
        try:
            file_extension = Path(file_path).suffix.lower()

            if file_extension == '.pdf':
                return TextExtractionService.extract_pdf(file_path)
            elif file_extension == '.docx':
                return TextExtractionService.extract_docx(file_path)
            elif file_extension == '.doc':
                return TextExtractionService.extract_doc(file_path)
            elif file_extension == '.pptx':
                return TextExtractionService.extract_pptx(file_path)
            elif file_extension in {
                '.txt', '.md', '.html', '.css', '.json', '.xml', '.yaml', '.yml'
            } or file_extension in TextExtractionService.CODE_EXTENSIONS:
                return TextExtractionService.extract_text(file_path)
            else:
                logger.warning(f"Unsupported file format: {file_extension}")
                return ""

        except Exception as e:
            logger.error(f"Error extracting text from file: {str(e)}")
            return ""

    @staticmethod
    def extract_pdf(file_path: str) -> str:
        """Extract text from PDF"""
        if not PyPDF2:
            logger.warning("PyPDF2 not installed")
            return ""

        try:
            text = ""
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                for page_num in range(len(pdf_reader.pages)):
                    page = pdf_reader.pages[page_num]
                    text += page.extract_text()
            return text
        except Exception as e:
            logger.error(f"Error extracting PDF: {str(e)}")
            return ""

    @staticmethod
    def extract_docx(file_path: str) -> str:
        """Extract text from DOCX"""
        if not Document:
            logger.warning("python-docx not installed")
            return ""

        try:
            doc = Document(file_path)
            text = ""
            for para in doc.paragraphs:
                text += para.text + "\n"
            
            # Also extract from tables
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        text += cell.text + " "
            return text
        except Exception as e:
            logger.error(f"Error extracting DOCX: {str(e)}")
            return ""

    @staticmethod
    def extract_doc(file_path: str) -> str:
        """Extract text from DOC (try as DOCX or plain text)"""
        # Try as DOCX first
        try:
            return TextExtractionService.extract_docx(file_path)
        except:
            # Try as plain text
            try:
                with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                    return f.read()
            except Exception as e:
                logger.error(f"Error extracting DOC: {str(e)}")
                return ""

    @staticmethod
    def extract_pptx(file_path: str) -> str:
        """Extract text from PPTX"""
        if not Presentation:
            logger.warning("python-pptx not installed")
            return ""

        try:
            prs = Presentation(file_path)
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            return text
        except Exception as e:
            logger.error(f"Error extracting PPTX: {str(e)}")
            return ""

    @staticmethod
    def extract_text(file_path: str) -> str:
        """Extract plain text from TXT file"""
        try:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                return f.read()
        except Exception as e:
            logger.error(f"Error extracting TXT: {str(e)}")
            return ""

    @staticmethod
    def extract_from_url(file_url: str) -> str:
        """
        Download file from URL and extract text
        
        Args:
            file_url (str): URL to the file
            
        Returns:
            str: Extracted text
        """
        import requests
        import tempfile

        try:
            response = requests.get(file_url, timeout=30)
            response.raise_for_status()

            parsed = urlparse(file_url)
            ext = Path(parsed.path).suffix.lower()
            if not ext:
                ext = TextExtractionService._extension_from_content_type(
                    response.headers.get("Content-Type", "")
                )
            if not ext:
                ext = TextExtractionService._extension_from_magic(response.content)
            if not ext:
                ext = ".txt"

            # Save to temp file
            with tempfile.NamedTemporaryFile(delete=False, suffix=ext) as tmp:
                tmp.write(response.content)
                tmp_path = tmp.name

            # Extract text
            text = TextExtractionService.extract_from_file(tmp_path)

            # Clean up
            if os.path.exists(tmp_path):
                os.remove(tmp_path)

            if text:
                return text

            # Fallback for unknown/plain content without supported extension.
            return response.content.decode('utf-8', errors='ignore')

        except Exception as e:
            logger.error(f"Error extracting from URL: {str(e)}")
            return ""

    @staticmethod
    def extract_from_bytes(file_bytes: bytes, file_name: str = "submission.txt") -> str:
        """
        Extract text from raw file bytes.

        Args:
            file_bytes (bytes): File content bytes
            file_name (str): Original filename to infer extension

        Returns:
            str: Extracted text
        """
        import tempfile

        if not file_bytes:
            return ""

        ext = Path(file_name or "submission.txt").suffix.lower()
        if not ext:
            ext = TextExtractionService._extension_from_magic(file_bytes) or ".txt"

        tmp_path = None
        try:
            with tempfile.NamedTemporaryFile(delete=False, suffix=ext) as tmp:
                tmp.write(file_bytes)
                tmp_path = tmp.name

            text = TextExtractionService.extract_from_file(tmp_path)
            if text:
                return text

            return file_bytes.decode('utf-8', errors='ignore')
        except Exception as e:
            logger.error(f"Error extracting from bytes: {str(e)}")
            return ""
        finally:
            if tmp_path and os.path.exists(tmp_path):
                os.remove(tmp_path)

    @staticmethod
    def is_code_file(file_name: str) -> bool:
        """Check whether a filename likely contains source code."""
        ext = Path(file_name or "").suffix.lower()
        return ext in TextExtractionService.CODE_EXTENSIONS
